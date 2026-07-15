import io
import ezdxf
from pptx import Presentation
from pptx.util import Inches, Pt

def generate_dxf_file(mapping_data):
    """
    Tạo cấu trúc các lớp riêng biệt trong CAD và vẽ hình học thực thể.
    """
    doc = ezdxf.new('R2010')
    msp = doc.modelspace()
    
    # Định dạng các layer theo tiêu chuẩn kỹ thuật thiết kế
    doc.layers.new(name='LAYER_BASE', dxfattribs={'color': 7})        # Lưới nền (Trắng/Đen)
    doc.layers.new(name='LAYER_MEP_DEVICES', dxfattribs={'color': 1}) # Thiết bị (Đỏ)
    doc.layers.new(name='LAYER_TEXT_LABELS', dxfattribs={'color': 3}) # Ghi chú (Xanh lá)
    
    # Vẽ hệ lưới 6x6 mét làm phông nền gốc
    for i in range(7):
        msp.add_line((i, 0), (i, 6), dxfattribs={'layer': 'LAYER_BASE'})
        msp.add_line((0, i), (6, i), dxfattribs={'layer': 'LAYER_BASE'})
        
    # Duyệt qua dữ liệu mapping để vẽ đè các điểm thiết bị
    for item in mapping_data:
        x = float(item.get('x', 0))
        y = float(item.get('y', 0))
        name = item.get('item_name', 'Device')
        
        # Thêm thực thể hình học Circle tại tọa độ chuẩn xác
        msp.add_circle((x, y), radius=0.1, dxfattribs={'layer': 'LAYER_MEP_DEVICES'})
        # Gắn nhãn Text thực thể cạnh vòng tròn
        msp.add_text(name[:6].upper(), dxfattribs={'layer': 'LAYER_TEXT_LABELS'}).set_placement((x + 0.15, y))
        
    # Ghi dữ liệu vào bộ nhớ đệm luồng byte để Streamlit cho phép download trực tiếp
    out_stream = io.StringIO()
    doc.write(out_stream)
    byte_stream = io.BytesIO(out_stream.getvalue().encode('utf-8'))
    return byte_stream.getvalue()

def generate_pptx_file(mapping_data):
    """
    Sinh tự động bộ slide báo cáo thi công gồm đúng 4 slide tiêu chuẩn bắt buộc.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Định dạng chuẩn tỉ lệ màn hình 16:9
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Sử dụng layout trống để tùy biến tọa độ phần tử
    
    # ---- Slide 1: Booth Location ----
    s1 = prs.slides.add_slide(blank_layout)
    t1 = s1.shapes.add_textbox(Inches(1), Inches(3), Inches(113), Inches(1))
    t1.text_frame.text = "SLIDE 1: BOOTH LOCATION - SƠ ĐỒ VỊ TRÍ GIAN HÀNG TRÊN MẶT BẰNG"
    
    # ---- Slide 2: Perspective View ----
    s2 = prs.slides.add_slide(blank_layout)
    t2 = s2.shapes.add_textbox(Inches(1), Inches(3), Inches(113), Inches(1))
    t2.text_frame.text = "SLIDE 2: PERSPECTIVE VIEW - PHỐI CẢNH MÔ HÌNH KHÔNG GIAN 3D"
    
    # ---- Slide 3: Booth Dimensions ----
    s3 = prs.slides.add_slide(blank_layout)
    t3 = s3.shapes.add_textbox(Inches(1), Inches(3), Inches(113), Inches(1))
    t3.text_frame.text = "SLIDE 3: BOOTH DIMENSIONS - BẢN VẼ KÍCH THƯỚC KHUNG HỆ LƯỚI GỐC"
    
    # ---- Slide 4: MEP Layout Chi Tiết Nâng Cao & Legend ----
    s4 = prs.slides.add_slide(blank_layout)
    
    # Tiêu đề Slide 4
    title_box = s4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = "SLIDE 4: MEP LAYOUT CHI TIẾT NÂNG CAO & BẢNG KÝ HIỆU"
    p.font.size = Pt(24)
    p.font.bold = True
    
    # Tự động hóa tích hợp bảng Ghi chú ký hiệu (Legend Table) ở góc phải slide
    rows = len(mapping_data) + 1
    cols = 3
    left = Inches(8.0)
    top = Inches(1.5)
    width = Inches(4.8)
    height = Inches(0.35 * rows)
    
    table_shape = s4.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Định dạng tiêu đề cột cho bảng
    table.cell(0, 0).text = "Thiết Bị Kỹ Thuật"
    table.cell(0, 1).text = "SL"
    table.cell(0, 2).text = "Tọa độ (X, Y)"
    
    # Điền tự động dữ liệu thiết bị thực tế khớp 100% với báo giá
    for idx, item in enumerate(mapping_data):
        table.cell(idx + 1, 0).text = str(item.get('item_name'))
        table.cell(idx + 1, 1).text = "1"
        table.cell(idx + 1, 2).text = f"({item.get('x')}, {item.get('y')})"
        
    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()
