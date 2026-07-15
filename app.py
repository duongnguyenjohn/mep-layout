"""
====================================================================================
 HỆ THỐNG TỰ ĐỘNG BÓC TÁCH & VẼ SƠ ĐỒ ĐIỆN GIAN HÀNG (Electrical Layout Auto-Generator)
 PHIÊN BẢN V6.0 — AI Smart Cropping (Tự cắt Top View) + Caching Component
====================================================================================
"""

import io
import os
import random
import re
import json
import unicodedata
import base64
import gc
import tempfile
from dataclasses import dataclass, field
from typing import Optional

import pandas as pd
import streamlit as st
import streamlit.components.v1 as components
from PIL import Image
from pypdf import PdfReader
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

try:
    from google import genai
except ImportError:
    genai = None

# ====================================================================================
# 0. KHỞI TẠO CUSTOM COMPONENT BẰNG CACHE (ĐẢM BẢO 100% KHÔNG LỖI LOAD)
# ====================================================================================
COMPONENT_HTML = """
<!DOCTYPE html>
<html lang="vi">
<head>
    <meta charset="UTF-8">
    <script src="https://cdn.tailwindcss.com"></script>
    <script src="https://html2canvas.hertzen.com/dist/html2canvas.min.js"></script>
    <script src="https://unpkg.com/streamlit-component-lib@1.3.0/dist/streamlit.js"></script>
    <style>
        body { font-family: 'Segoe UI', Arial, sans-serif; margin: 0; padding: 10px; background: transparent; }
        .map-container { position: relative; width: 800px; height: 800px; background-size: 100% 100%; border: 2px solid #005088; margin: 0 auto; overflow: hidden; background-color: #f0f0f0; }
        .grid-overlay { position: absolute; top: 0; left: 0; right: 0; bottom: 0; background-image: linear-gradient(rgba(0, 80, 136, 0.4) 1px, transparent 1px), linear-gradient(90deg, rgba(0, 80, 136, 0.4) 1px, transparent 1px); background-size: 133.33px 133.33px; pointer-events: none; z-index: 1; }
        .icon-node { position: absolute; z-index: 10; cursor: grab; transform: translate(-50%, -50%); display: flex; align-items: center; justify-content: center; width: 30px; height: 30px; border-radius: 4px; color: white; font-size: 14px; font-weight: bold; box-shadow: 0 2px 5px rgba(0,0,0,0.5); user-select: none; }
        .icon-node:active { cursor: grabbing; z-index: 999; }
        .toolbar { width: 800px; margin: 10px auto; padding: 10px; background: #fff; border: 1px solid #ccc; border-radius: 8px; display: flex; gap: 10px; align-items: center; justify-content: space-between; }
        .btn { padding: 8px 16px; background: #005088; color: white; border: none; border-radius: 4px; cursor: pointer; font-weight: bold; }
        .btn:hover { background: #003057; }
        .btn-success { background: #10b981; }
        .btn-success:hover { background: #059669; }
        .delete-hint { font-size: 12px; color: #666; }
    </style>
</head>
<body>
    <div class="toolbar">
        <div><span class="delete-hint">💡 Kéo thả để di chuyển. Click chuột phải vào Icon để XÓA.</span></div>
        <div><button class="btn btn-success" id="btn-save" onclick="saveAndSend()">📸 Chốt Bản Đồ & Lưu</button></div>
    </div>
    <div class="map-container" id="map-area">
        <div class="grid-overlay"></div>
    </div>

    <script>
        let isDragging = false;
        let currentDrag = null;
        let itemsData = [];

        function initComponent() { Streamlit.setFrameHeight(900); }

        function onDataFromPython(event) {
            if (event.data.type !== "streamlit:render") return;
            const args = event.data.args;
            const mapArea = document.getElementById("map-area");
            if (args.bg_base64) mapArea.style.backgroundImage = `url(${args.bg_base64})`;
            if (itemsData.length === 0 && args.items) {
                itemsData = args.items;
                renderItems();
            }
        }

        function renderItems() {
            const mapArea = document.getElementById("map-area");
            mapArea.querySelectorAll('.icon-node').forEach(e => e.remove());
            itemsData.forEach((item, index) => {
                const node = document.createElement("div");
                node.className = "icon-node";
                node.innerHTML = item.icon_val;
                node.style.backgroundColor = item.color;
                node.style.left = `${(item.x / 6) * 100}%`;
                node.style.top = `${(item.y / 6) * 100}%`;
                node.dataset.index = index;

                node.addEventListener("mousedown", (e) => { if (e.button === 0) { isDragging = true; currentDrag = node; } });
                node.addEventListener("contextmenu", (e) => { e.preventDefault(); if(confirm("Xóa thiết bị này?")) node.remove(); });
                mapArea.appendChild(node);
            });
        }

        document.addEventListener("mousemove", (e) => {
            if (isDragging && currentDrag) {
                const rect = document.getElementById("map-area").getBoundingClientRect();
                let x = Math.max(0, Math.min(e.clientX - rect.left, rect.width));
                let y = Math.max(0, Math.min(e.clientY - rect.top, rect.height));
                currentDrag.style.left = `${(x / rect.width) * 100}%`;
                currentDrag.style.top = `${(y / rect.height) * 100}%`;
            }
        });

        document.addEventListener("mouseup", () => { isDragging = false; currentDrag = null; });

        function saveAndSend() {
            const mapArea = document.getElementById("map-area");
            const btn = document.getElementById("btn-save");
            btn.innerText = "⏳ Đang chụp màn hình...";
            
            let finalCounts = {};
            document.querySelectorAll('.icon-node').forEach(node => {
                let key = itemsData[node.dataset.index].key;
                finalCounts[key] = (finalCounts[key] || 0) + 1;
            });

            html2canvas(mapArea, { useCORS: true, scale: 2 }).then(canvas => {
                Streamlit.setComponentValue({ status: "approved", final_image_b64: canvas.toDataURL("image/png"), final_counts: finalCounts });
                btn.innerText = "✅ Đã gửi! Vui lòng bấm Xuất File trên web.";
                btn.classList.remove('btn-success');
            });
        }

        Streamlit.events.addEventListener(Streamlit.RENDER_EVENT, onDataFromPython);
        Streamlit.setComponentReady();
        initComponent();
    </script>
</body>
</html>
"""

@st.cache_resource
def get_interactive_map_component():
    """Tạo component 1 lần duy nhất trong RAM để tránh mọi lỗi load file"""
    component_dir = os.path.join(tempfile.gettempdir(), "mep_map_v6")
    os.makedirs(component_dir, exist_ok=True)
    index_path = os.path.join(component_dir, "index.html")
    with open(index_path, "w", encoding="utf-8") as f:
        f.write(COMPONENT_HTML)
    return components.declare_component("interactive_map", path=component_dir)

interactive_map_component = get_interactive_map_component()

# ====================================================================================
# 1. CẤU HÌNH MẶC ĐỊNH
# ====================================================================================
AI_MODEL_DEFAULT = "gemini-1.5-pro"
STOPWORDS_VI = {"và", "cho", "của", "tại"}
SLIDE_TYPES = ["Booth Location", "Perspective View", "Booth Dimensions"]

DEFAULT_EQUIPMENT_CONFIG = [
    {"Tên thiết bị": "Đèn Floodlight 50W", "Biểu tượng vẽ": "▲", "Mã màu Hex": "#FFCD00", "Vị trí ưu tiên": "Hệ trần biên", "Công suất": "50W", "Số lượng": 30},
    {"Tên thiết bị": "Ổ cắm 5A/220V", "Biểu tượng vẽ": "●", "Mã màu Hex": "#D62728", "Vị trí ưu tiên": "Bàn tư vấn", "Công suất": "220V", "Số lượng": 3},
]

# ====================================================================================
# 2. DATA STRUCTURES & HELPER
# ====================================================================================
@dataclass
class EquipmentRow:
    key: str
    label: str
    icon_value: str
    color_hex: str
    zone_raw: str
    power: str
    quantity: int

@dataclass
class PipelineResult:
    equipment: list = field(default_factory=list)
    booth_location_img: Optional[Image.Image] = None
    perspective_img: Optional[Image.Image] = None
    dimensions_img: Optional[Image.Image] = None
    coordinates: dict = field(default_factory=dict)

def slugify(text: str, idx: int) -> str:
    try:
        norm = unicodedata.normalize("NFKD", text or "").encode("ascii", "ignore").decode("ascii")
        norm = re.sub(r"[^a-zA-Z0-9]+", "_", norm).strip("_").lower()
        return f"{norm}_{idx}"
    except: return f"item_{idx}"

def read_quotation_as_text(file_bytes: bytes, filename: str) -> str:
    ext = os.path.splitext(filename.lower())[1]
    try:
        if ext == ".pdf":
            reader = PdfReader(io.BytesIO(file_bytes))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext in (".xlsx", ".xls"):
            sheets = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
            return "\n".join(df.to_string(index=False) for df in sheets.values())
        else: return ""
    except: return ""

def extract_salient_tokens(label: str) -> list:
    tokens = re.split(r"[\s/()\-,]+", (label or "").lower())
    return [t for t in tokens if len(t) >= 2 and t not in STOPWORDS_VI]

def regex_extract_quantities(text: str, labels: list) -> dict:
    found = {}
    lines = text.lower().splitlines()
    qty_patterns = [r"(?:sl|qty|số lượng)\s*[:.\-]?\s*(\d+)", r"x\s*(\d+)\b", r"\b(\d+)\s*(?:bộ|cái|chiếc|nos|no)\b"]
    for label in labels:
        tokens = extract_salient_tokens(label)
        if not tokens: continue
        best_qty = 0
        for line in lines:
            if any(tok in line for tok in tokens):
                for pat in qty_patterns:
                    m = re.search(pat, line)
                    if m: best_qty = max(best_qty, int(m.group(1)))
        if best_qty > 0: found[label] = best_qty
    return found

# ====================================================================================
# 3. AI PROVIDER (SMART CROPPING - CẮT TOP VIEW)
# ====================================================================================
class AIProvider:
    def __init__(self, api_key: str, model: str):
        if genai is None:
            raise RuntimeError("Chưa cài đặt google-genai.")
        self.client = genai.Client(api_key=api_key)
        self.model = model

    def analyze_slide_image(self, image: Image.Image) -> dict:
        system_prompt = """Bạn là chuyên gia phân tích bản vẽ PDF. 
Trang PDF này có thể chứa một hoặc nhiều góc nhìn (Top view, Front view, Right view...) ghép chung.

Tìm các từ khóa như "Booth Dimensions", "Top view" trên ảnh để định vị.
Trả về ĐÚNG định dạng JSON (không có markdown):
{
    "label": "Booth Location" | "Perspective View" | "Booth Dimensions" | "Other",
    "crop_box": [ymin, xmin, ymax, xmax]
}

Quy tắc:
1. Nếu trang LÀ mặt bằng tổng thể toàn khu, trả về label: "Booth Location".
2. Nếu trang LÀ phối cảnh 3D, trả về label: "Perspective View".
3. Nếu trang có tiêu đề "Booth Dimensions" HOẶC chứa bản vẽ "Top view" (Mặt bằng 2D nhìn từ trên xuống có lưới tọa độ):
   - Trả về label: "Booth Dimensions".
   - Bạn PHẢI khoanh vùng (Bounding Box) cắt RIÊNG phần bản vẽ "Top view" đó ra. Bỏ qua các góc nhìn Front/Left/Right.
   - Cắt rộng ra một chút để lấy ĐẦY ĐỦ các chữ số tọa độ (0, 1, 2, 3...) ở viền ngoài lưới.
   - Tọa độ `crop_box` là mảng 4 số nguyên từ 0 đến 1000 (tỷ lệ 0%-100% kích thước ảnh). Ví dụ: [100, 50, 900, 450].
4. Nếu không thuộc các loại trên, để crop_box là null."""
        
        try:
            response = self.client.models.generate_content(
                model=self.model,
                contents=[system_prompt, image]
            )
            text = response.text.strip()
            # Bắt JSON bằng Regex để tránh lỗi định dạng markdown của Gemini
            match = re.search(r'\{.*\}', text, re.DOTALL)
            if match:
                return json.loads(match.group(0))
            return {"label": "Other", "crop_box": None}
        except Exception as e:
            print("AI Vision Error:", e)
            return {"label": "Other", "crop_box": None}

# ====================================================================================
# 4. RUN AI PROCESSING
# ====================================================================================
def run_ai_processing(quotation_file, review_file, table_rows: list, ai: AIProvider) -> PipelineResult:
    result = PipelineResult()
    
    quotation_text = read_quotation_as_text(quotation_file.read(), quotation_file.name)
    labels = [str(r.get("Tên thiết bị", "")).strip() for r in table_rows if str(r.get("Tên thiết bị", "")).strip()]
    regex_found = regex_extract_quantities(quotation_text, labels) if quotation_text else {}
    
    for idx, row in enumerate(table_rows):
        label = str(row["Tên thiết bị"]).strip()
        synced_qty = regex_found.get(label, int(row["Số lượng"])) 
        result.equipment.append(EquipmentRow(
            key=slugify(label, idx), label=label, icon_value=row["Biểu tượng vẽ"],
            color_hex=row["Mã màu Hex"], zone_raw=row["Vị trí ưu tiên"], 
            power=row["Công suất"], quantity=int(synced_qty)
        ))
    
    images = convert_from_bytes(review_file.read(), dpi=100)
    classified = {t: None for t in SLIDE_TYPES}
    
    for img in images:
        data = ai.analyze_slide_image(img)
        label = data.get("label", "Other")
        
        if label == "Booth Dimensions" and classified["Booth Dimensions"] is None:
            crop_box = data.get("crop_box")
            if crop_box and len(crop_box) == 4:
                try:
                    # AI đã chỉ ra tọa độ, ta dùng Pillow để "Cắt" riêng Top View ra
                    ymin, xmin, ymax, xmax = crop_box
                    w, h = img.size
                    left, top = (xmin / 1000.0) * w, (ymin / 1000.0) * h
                    right, bottom = (xmax / 1000.0) * w, (ymax / 1000.0) * h
                    classified["Booth Dimensions"] = img.crop((left, top, right, bottom))
                except Exception:
                    classified["Booth Dimensions"] = img
            else:
                classified["Booth Dimensions"] = img
                
        elif label == "Booth Location" and classified["Booth Location"] is None:
            classified["Booth Location"] = img
        elif label == "Perspective View" and classified["Perspective View"] is None:
            classified["Perspective View"] = img
            
    result.booth_location_img = classified["Booth Location"]
    result.perspective_img = classified["Perspective View"]
    result.dimensions_img = classified["Booth Dimensions"]
    
    coords = {}
    for it in result.equipment:
        coords[it.key] = [[round(random.uniform(1.0, 5.0), 2), round(random.uniform(1.0, 5.0), 2)] for _ in range(it.quantity)]
    result.coordinates = coords
    
    del images
    gc.collect()
    
    return result

# ====================================================================================
# 5. HÀM TẠO PPTX
# ====================================================================================
def export_final_pptx(result: PipelineResult, final_b64_image: str, final_counts: dict) -> bytes:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_img(slide, img_obj):
        img_w, img_h = img_obj.size
        slide_ratio = prs.slide_width / prs.slide_height
        pic_w, pic_h = (prs.slide_width, int(prs.slide_width / (img_w/img_h))) if (img_w/img_h) > slide_ratio else (int(prs.slide_height * (img_w/img_h)), prs.slide_height)
        left, top = int((prs.slide_width - pic_w)/2), int((prs.slide_height - pic_h)/2)
        buf = io.BytesIO()
        img_obj.convert("RGB").save(buf, format="PNG")
        slide.shapes.add_picture(buf, left, top, width=pic_w, height=pic_h)

    for img in [result.booth_location_img, result.perspective_img, result.dimensions_img]:
        s = prs.slides.add_slide(blank)
        if img: add_img(s, img)

    s4 = prs.slides.add_slide(blank)
    if final_b64_image:
        header, encoded = final_b64_image.split(",", 1)
        final_img = Image.open(io.BytesIO(base64.b64decode(encoded)))
        add_img(s4, final_img)

        rows = []
        for eq in result.equipment:
            qty = final_counts.get(eq.key, 0)
            if qty > 0:
                rows.append((eq.label, eq.icon_value, eq.color_hex, str(qty)))
        
        if rows:
            table_shape = s4.shapes.add_table(len(rows)+1, 4, prs.slide_width - Inches(4.5), Inches(0.25), Inches(4.4), Inches(0.3 * (len(rows)+1)))
            for c, t in enumerate(["Thiết bị", "Ký hiệu", "Màu", "SL"]):
                table_shape.table.cell(0, c).text = t
            for r, (lbl, sym, col, qt) in enumerate(rows, 1):
                table_shape.table.cell(r, 0).text = lbl
                table_shape.table.cell(r, 1).text = sym
                table_shape.table.cell(r, 3).text = qt
                table_shape.table.cell(r, 2).fill.solid()
                table_shape.table.cell(r, 2).fill.fore_color.rgb = RGBColor(*[int(col.strip('#')[i:i+2], 16) for i in (0, 2, 4)])

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

# ====================================================================================
# 6. GIAO DIỆN CHÍNH
# ====================================================================================
def main():
    st.set_page_config(page_title="MEP Layout Auto", page_icon="⚡", layout="wide")
    
    if "step" not in st.session_state:
        st.session_state.step = 1
        st.session_state.result = None
        st.session_state.equipment_table = pd.DataFrame(DEFAULT_EQUIPMENT_CONFIG)

    st.title("⚡ MEP Layout Auto-Generator V6.0 (Smart Crop)")
    
    if interactive_map_component is None:
        st.error("❌ Lỗi hệ thống: Không thể khởi tạo Component Web. Vui lòng liên hệ Admin.")
        return

    if st.session_state.step == 1:
        st.info("Bước 1: Cấu hình bảng Legend và Upload file thiết kế (Quotation + 3D Review).")
        c1, c2 = st.columns([1, 1])
        with c1:
            edited_df = st.data_editor(st.session_state.equipment_table, num_rows="dynamic", use_container_width=True)
            st.session_state.equipment_table = edited_df
        with c2:
            st.sidebar.text_input("Google Gemini API Key", key="api_key", type="password")
            st.sidebar.selectbox("Model", ["gemini-1.5-pro", "gemini-2.5-flash"], key="model_choice")
            
            q_file = st.file_uploader("File Báo Giá (PDF/Excel)")
            r_file = st.file_uploader("File 3D Review (PDF)")
            
            if st.button("🚀 Xử Lý AI Khởi Tạo Bản Đồ"):
                if not st.session_state.api_key: 
                    st.error("⚠️ Vui lòng nhập API Key!")
                elif not q_file or not r_file: 
                    st.error("⚠️ Vui lòng upload đầy đủ 2 file!")
                else:
                    try:
                        ai = AIProvider(st.session_state.api_key, st.session_state.model_choice)
                        with st.spinner("AI đang tìm và tự động cắt bản vẽ Top View..."):
                            st.session_state.result = run_ai_processing(q_file, r_file, edited_df.to_dict("records"), ai)
                            st.session_state.step = 2
                            st.rerun()
                    except Exception as e:
                        st.error(f"❌ Có lỗi xảy ra trong quá trình xử lý AI: {str(e)}")

    elif st.session_state.step == 2:
        st.success("Bước 2: Hệ thống đã tự động cắt bản vẽ Top View! Hãy kéo thả Icon cho đúng vị trí, rồi bấm 'Chốt Bản Đồ & Lưu'.")
        
        if st.button("⬅️ Quay lại Bước 1"):
            st.session_state.step = 1
            st.rerun()

        result = st.session_state.result
        bg_b64 = ""
        if result.dimensions_img:
            buf = io.BytesIO()
            result.dimensions_img.convert("RGB").save(buf, format="PNG")
            bg_b64 = f"data:image/png;base64,{base64.b64encode(buf.getvalue()).decode('utf-8')}"
        else:
            st.warning("⚠️ AI không tìm thấy hoặc cắt lỗi bản vẽ Top View chuẩn 2D trong file PDF.")

        js_items = []
        for eq in result.equipment:
            pts = result.coordinates.get(eq.key, [])
            for p in pts:
                js_items.append({"key": eq.key, "icon_val": eq.icon_value, "color": eq.color_hex, "x": p[0], "y": p[1]})

        component_value = interactive_map_component(
            bg_base64=bg_b64, 
            items=js_items, 
            key="map_interact"
        )

        if component_value and component_value.get("status") == "approved":
            st.divider()
            st.header("📥 BƯỚC CUỐI: Duyệt & Xuất File PPTX")
            
            final_b64 = component_value.get("final_image_b64")
            final_counts = component_value.get("final_counts", {})

            pptx_bytes = export_final_pptx(result, final_b64, final_counts)
            
            st.download_button(
                label="🎉 TẢI XUỐNG FILE MEP LAYOUT (.PPTX)",
                data=pptx_bytes,
                file_name="MEP_Layout_Final.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary"
            )

if __name__ == "__main__":
    main()
